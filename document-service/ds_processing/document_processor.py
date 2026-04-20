"""
Document Processor - Handles document text extraction and chunking
"""

import asyncio
import email
import json
import logging
import os
import tempfile
import zipfile
from pathlib import Path
from typing import List, Dict, Any
import time
import re

import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import ocrmypdf
from PIL import Image
from langdetect import detect
import textstat

from ds_config import settings
from ds_models.api_models import Chunk, QualityMetrics, ProcessingResult, Entity
from ds_processing.ocr_service import OCRService
from ds_processing.ocr_in_progress import add as ocr_in_progress_add, remove as ocr_in_progress_remove

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Processes documents and extracts text with quality assessment"""
    
    _instance = None
    _initialized = False
    _initialization_lock = asyncio.Lock()
    
    def __new__(cls):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
            cls._instance.document_service_client = None
            cls._instance.ocr_service = None
        return cls._instance
    
    def __init__(self):
        # __init__ is called every time, but we only want to initialize once
        # The actual initialization happens in the initialize() method
        pass
    
    async def initialize(self):
        """Initialize NLP models and OCR service (singleton pattern)"""
        async with self._initialization_lock:
            if self._initialized:
                logger.debug("🔄 DocumentProcessor already initialized, skipping")
                return
                
            try:
                logger.info("Initializing DocumentProcessor singleton...")

                # Standalone document-service: NER client is injected by pipeline
                # (DocumentProcessingPipeline sets document_service_client to _InternalNERClient).
                self.document_service_client = None

                # Initialize OCR service
                try:
                    self.ocr_service = OCRService()
                    await self.ocr_service.initialize()
                    logger.info("✅ OCR service initialized in document processor")
                except Exception as e:
                    logger.warning(f"⚠️  OCR service initialization failed: {e}")
                    self.ocr_service = None
                
                self._initialized = True
                logger.info("✅ DocumentProcessor singleton initialized successfully")
                
            except Exception as e:
                logger.error(f"❌ DocumentProcessor initialization failed: {e}")
                self._initialized = True  # Continue without advanced features
    
    @property
    def initialized(self) -> bool:
        """Check if the DocumentProcessor is initialized"""
        return self._initialized
    
    @classmethod
    def get_instance(cls):
        """Get the singleton instance"""
        return cls()
    
    @classmethod
    def reset_instance(cls):
        """Reset the singleton instance (mainly for testing)"""
        cls._instance = None
        cls._initialized = False
    
    def get_status(self) -> Dict[str, Any]:
        """Get the current status of the DocumentProcessor"""
        return {
            "initialized": self._initialized,
            "document_service_available": self.document_service_client is not None and getattr(
                self.document_service_client, "_initialized", False
            ),
            "spacy_model": None,
            "ocr_service_loaded": self.ocr_service is not None,
            "instance_id": id(self),
        }

    async def process_document(self, file_path: str, doc_type: str, document_id: str = None) -> ProcessingResult:
        """Process a document and return chunks with quality metrics
        
        Args:
            file_path: Path to the document file
            doc_type: Type of document (pdf, docx, etc.)
            document_id: UUID of the document (if None, derived from filename for backward compatibility)
        """
        start_time = time.time()
        
        # **ROOSEVELT FIX**: Use provided document_id instead of deriving from filename
        if document_id is None:
            document_id = Path(file_path).stem.split('_')[0]
            logger.warning(f"⚠️ No document_id provided, deriving from filename: {document_id}")
        
        try:
            logger.info(f"🔄 Processing {doc_type} document: {file_path} (doc_id: {document_id})")
            
            # ROOSEVELT DOCTRINE: Org Mode files use structured mechanisms, not vectorization!
            if doc_type == 'org':
                logger.info("Skipping vectorization for Org Mode file (structured data, not prose)")
                logger.info(f"📋 Org file will be handled by llm-orchestrator (OrgInboxAgent and OrgProjectAgent migrated)")
                
                # Return empty processing result - file is stored but not vectorized
                result = ProcessingResult(
                    document_id=document_id,
                    chunks=[],  # No chunks - structural queries only
                    entities=[],  # No entities - org-specific metadata used instead
                    quality_metrics=QualityMetrics(
                        overall_score=1.0,
                        ocr_confidence=1.0,
                        language_confidence=1.0,  # Default for structured data
                        vocabulary_score=1.0,  # Default for structured data
                        pattern_score=1.0  # Default for structured data
                    ),
                    processing_time=time.time() - start_time
                )
                logger.info(f"✅ Org file registered for structured access (no vectorization)")
                return result
            
            # Images are stored but NOT vectorized!
            if doc_type == 'image':
                logger.info(f"⏭️ Skipping vectorization for image file (binary data, not text)")
                logger.info(f"📷 Image stored for reference but not embedded")
                
                # Return empty processing result - file is stored but not vectorized
                result = ProcessingResult(
                    document_id=document_id,
                    chunks=[],  # No chunks - images not vectorized
                    entities=[],  # No entities
                    quality_metrics=QualityMetrics(
                        overall_score=1.0,
                        ocr_confidence=1.0,
                        language_confidence=1.0,
                        vocabulary_score=1.0,
                        pattern_score=1.0
                    ),
                    processing_time=time.time() - start_time
                )
                logger.info(f"✅ Image file registered for storage (no vectorization)")
                return result

            if doc_type in ('mp4', 'mkv', 'avi', 'mov', 'webm', 'audio'):
                logger.info(
                    "Skipping vectorization for binary media file (%s)",
                    doc_type,
                )
                return ProcessingResult(
                    document_id=document_id,
                    chunks=[],
                    entities=[],
                    quality_metrics=QualityMetrics(
                        overall_score=1.0,
                        ocr_confidence=1.0,
                        language_confidence=1.0,
                        vocabulary_score=1.0,
                        pattern_score=1.0,
                    ),
                    processing_time=time.time() - start_time,
                )
            
            # Extract text based on document type
            page_boundaries: list[tuple[int, int]] = []
            doc_type_hint: str | None = None
            if doc_type == 'pdf':
                text, ocr_confidence, page_boundaries, doc_type_hint = await self._process_pdf(file_path, document_id)
            elif doc_type == 'docx':
                text, ocr_confidence = await self._process_docx(file_path), 1.0
            elif doc_type == 'pptx':
                text, ocr_confidence = await self._process_pptx(file_path), 1.0
            elif doc_type == 'epub':
                text, ocr_confidence = await self._process_epub(file_path), 1.0
            elif doc_type == 'txt':
                text, ocr_confidence = await self._process_text(file_path), 1.0
            elif doc_type == 'md':
                text, ocr_confidence = await self._process_text(file_path), 1.0
            elif doc_type == 'html':
                text, ocr_confidence = await self._process_html(file_path), 1.0
            elif doc_type == 'eml':
                text, ocr_confidence = await self._process_eml(file_path), 1.0
            elif doc_type == 'zip':
                text, ocr_confidence = await self._process_zip(file_path), 1.0
            elif doc_type == 'srt':
                text, ocr_confidence = await self._process_subtitle_file(file_path, 'srt'), 1.0
            elif doc_type == 'vtt':
                text, ocr_confidence = await self._process_subtitle_file(file_path, 'vtt'), 1.0
            elif doc_type == 'image_sidecar':
                text, ocr_confidence = await self._process_image_sidecar(file_path), 1.0
            else:
                raise ValueError(f"Unsupported document type: {doc_type}")
            
            # Assess text quality
            quality_metrics = await self._assess_quality(text, ocr_confidence)
            
            # Chunk the text (pass page boundaries for PDF so chunks get page_start/page_end)
            chunks = await self._chunk_text(text, file_path, document_id, page_boundaries=page_boundaries, doc_type_hint=doc_type_hint)
            
            # Extract entities from the text
            entities = await self._extract_entities(text, chunks)
            
            # Create processing result
            result = ProcessingResult(
                document_id=document_id,  # **ROOSEVELT FIX**: Use provided UUID, not filename
                chunks=chunks,
                entities=entities,
                quality_metrics=quality_metrics,
                processing_time=time.time() - start_time
            )
            
            logger.info(f"✅ Document processed: {len(chunks)} chunks, quality: {quality_metrics.overall_score:.2f}")
            return result
            
        except Exception as e:
            logger.error(f"❌ Document processing failed: {e}")
            raise
    
    def _detect_newspaper_pdf(self, file_path: str) -> bool:
        """Heuristic: True if PDF looks like newspaper/magazine (multi-column or headline-heavy)."""
        try:
            with pdfplumber.open(file_path) as pdf:
                pages_to_check = min(3, len(pdf.pages))
                if pages_to_check < 1:
                    return False
                multi_column_count = 0
                headline_like_count = 0
                for i in range(pages_to_check):
                    page = pdf.pages[i]
                    w = page.width
                    words = page.extract_words() or []
                    if len(words) > 20 and w > 0:
                        x_positions = [float(wd.get("x0", 0)) for wd in words]
                        mid = w / 2
                        left = sum(1 for x in x_positions if x < mid)
                        right = len(x_positions) - left
                        if left > 5 and right > 5:
                            multi_column_count += 1
                    raw = page.extract_text() or ""
                    lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
                    for line in lines:
                        if 10 <= len(line) <= 80 and (line.isupper() or (line[0].isupper() and not line.endswith(".") and len(line.split()) <= 12)):
                            headline_like_count += 1
                            break
                if multi_column_count >= 1:
                    return True
                if headline_like_count >= 2 and pages_to_check >= 2:
                    return True
                return False
        except Exception:
            return False

    def _extract_pdf_text_column_order(self, file_path: str) -> tuple[str, list[tuple[int, int]]]:
        """Extract PDF text in column order (left then right half per page) to avoid cross-column merge."""
        text = ""
        page_boundaries: list[tuple[int, int]] = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    w = page.width
                    h = page.height
                    if w and h:
                        left = page.crop((0, 0, w / 2, h))
                        right = page.crop((w / 2, 0, w, h))
                        left_text = (left.extract_text() or "").strip()
                        right_text = (right.extract_text() or "").strip()
                        page_text = left_text + "\n" + right_text if (left_text and right_text) else (left_text or right_text)
                    else:
                        page_text = page.extract_text() or ""
                    if page_text:
                        page_boundaries.append((i + 1, len(text)))
                        text += page_text + "\n"
        except Exception:
            pass
        return text, page_boundaries

    async def _process_pdf(self, file_path: str, document_id: str | None = None) -> tuple[str, float, list[tuple[int, int]], str | None]:
        """Process PDF document with automated fallback to OCR.
        Returns (text, ocr_confidence, page_boundaries, doc_type_hint).
        doc_type_hint is 'newspaper' when layout suggests newspaper/magazine, else None.
        """
        doc_id = document_id or ""
        text = ""
        ocr_confidence = 1.0
        page_boundaries: list[tuple[int, int]] = []
        doc_type_hint: str | None = None

        try:
            # Try standard text extraction (page-by-page for boundaries)
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        page_boundaries.append((i + 1, len(text)))
                        text += page_text + "\n"

            # If no text extracted, try pdfplumber (better for some layouts)
            if not text.strip():
                page_boundaries = []
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            page_boundaries.append((i + 1, len(text)))
                            text += page_text + "\n"

            # If still no text, the PDF is likely a scan - trigger OCR fallback
            if not text.strip():
                logger.info(f"No text found in {doc_id or file_path}, triggering OCR fallback")
                text, ocr_confidence, page_boundaries = await self._ocr_pdf(file_path)

            # Newspaper/magazine: re-extract in column order and hint chunker
            if text.strip() and self._detect_newspaper_pdf(file_path):
                col_text, col_boundaries = self._extract_pdf_text_column_order(file_path)
                if col_text.strip():
                    text = col_text
                    page_boundaries = col_boundaries
                    doc_type_hint = "newspaper"

        except Exception as e:
            logger.warning(f"PDF text extraction failed, trying OCR: {e}")
            text, ocr_confidence, page_boundaries = await self._ocr_pdf(file_path)

        return text, ocr_confidence, page_boundaries, doc_type_hint
    
    async def _ocr_pdf(self, file_path: str) -> tuple[str, float, list[tuple[int, int]]]:
        """Perform OCR on PDF. Registers path so file watcher skips re-processing during OCR.
        Returns (text, confidence, page_boundaries) with (page_num_1based, char_start) per page.
        """
        ocr_in_progress_add(file_path)
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                temp_path = temp_file.name

            ocrmypdf.ocr(
                file_path,
                temp_path,
                language='+'.join(settings.OCR_LANGUAGES),
                force_ocr=True,
                deskew=True,
                clean=True,
                optimize=1,
                output_type='pdf',
            )

            text = ""
            page_boundaries: list[tuple[int, int]] = []
            with pdfplumber.open(temp_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        page_boundaries.append((i + 1, len(text)))
                        text += page_text + "\n"

            if os.path.exists(temp_path):
                os.unlink(temp_path)

            confidence = min(1.0, len([c for c in text if c.isalnum()]) / max(1, len(text)) * 2)
            return text, confidence, page_boundaries

        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return "", 0.0, []
        finally:
            ocr_in_progress_remove(file_path)
    
    async def _process_docx(self, file_path: str) -> str:
        """Process DOCX document"""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"❌ DOCX processing failed: {e}")
            return ""

    async def _process_pptx(self, file_path: str) -> str:
        """Process PPTX document: extract text from slides, tables, and speaker notes."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = (para.text or "").strip()
                            if text:
                                slide_text.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                (cell.text or "").strip() for cell in row.cells
                            )
                            if row_text.strip():
                                slide_text.append(row_text)
                if slide_text:
                    parts.append(f"## Slide {slide_num}\n" + "\n".join(slide_text))
                try:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                        if notes:
                            parts.append(f"*Speaker notes:* {notes}")
                except Exception:
                    pass
            return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logger.error(f"PPTX processing failed: {e}")
            return ""
    
    async def _process_epub(self, file_path: str) -> str:
        """Process EPUB document"""
        try:
            book = epub.read_epub(file_path)
            text = ""
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_body_content(), 'html.parser')
                    text += soup.get_text() + "\n"
            
            return text
        except Exception as e:
            logger.error(f"❌ EPUB processing failed: {e}")
            return ""
    
    async def _process_text(self, file_path: str) -> str:
        """Process plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            logger.error(f"❌ Could not decode text file: {file_path}")
            return ""
        except Exception as e:
            logger.error(f"❌ Text processing failed: {e}")
            return ""

    async def _process_image_sidecar(self, file_path: str) -> str:
        """
        Extract searchable plain text from an image metadata JSON sidecar (*.metadata.json).
        Mirrors backend image_sidecar_service searchable_parts assembly.
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                metadata_data = json.load(f)
        except json.JSONDecodeError as e:
            logger.error("Image sidecar JSON invalid: %s: %s", file_path, e)
            raise ValueError(f"Invalid image sidecar JSON: {e}") from e
        except OSError as e:
            logger.error("Image sidecar read failed: %s: %s", file_path, e)
            raise

        title = (metadata_data.get("title") or "").strip()
        content = (
            metadata_data.get("content") or metadata_data.get("transcript") or ""
        ).strip()
        author = (metadata_data.get("author") or "").strip()
        date = (metadata_data.get("date") or "").strip()
        series = (metadata_data.get("series") or "").strip()
        tags = metadata_data.get("tags", []) or []
        location = (metadata_data.get("location") or "").strip()
        event = (metadata_data.get("event") or "").strip()
        medium = (metadata_data.get("medium") or "").strip()
        dimensions = (metadata_data.get("dimensions") or "").strip()
        body_part = (metadata_data.get("body_part") or "").strip()
        modality = (metadata_data.get("modality") or "").strip()
        map_type = (metadata_data.get("map_type") or "").strip()
        coordinates = (metadata_data.get("coordinates") or "").strip()
        application = (metadata_data.get("application") or "").strip()
        platform = (metadata_data.get("platform") or "").strip()

        if not title:
            raise ValueError("Image sidecar missing required field: title")
        if "content" not in metadata_data and "transcript" not in metadata_data:
            raise ValueError(
                "Image sidecar missing required field: content or transcript (key must exist)"
            )

        searchable_parts: list[str] = [title]
        if content:
            searchable_parts.append(content)
        if author:
            searchable_parts.append(author)
        if series:
            searchable_parts.append(series)
        if tags:
            searchable_parts.extend(tags)
        if date:
            searchable_parts.append(date)
        if location:
            searchable_parts.append(f"Location: {location}")
        if event:
            searchable_parts.append(f"Event: {event}")
        if medium:
            searchable_parts.append(f"Medium: {medium}")
        if dimensions:
            searchable_parts.append(f"Dimensions: {dimensions}")
        if body_part:
            searchable_parts.append(f"Body part: {body_part}")
        if modality:
            searchable_parts.append(f"Modality: {modality}")
        if map_type:
            searchable_parts.append(f"Map type: {map_type}")
        if coordinates:
            searchable_parts.append(f"Coordinates: {coordinates}")
        if application:
            searchable_parts.append(f"Application: {application}")
        if platform:
            searchable_parts.append(f"Platform: {platform}")

        faces_data = metadata_data.get("faces") or []
        if isinstance(faces_data, list):
            for face in faces_data:
                if isinstance(face, dict):
                    name = (
                        face.get("identity_name") or face.get("suggested_identity") or ""
                    ).strip()
                    if name:
                        searchable_parts.append(name)

        objects_data = metadata_data.get("objects") or []
        if isinstance(objects_data, list):
            for obj in objects_data:
                if isinstance(obj, dict):
                    label = (obj.get("user_tag") or obj.get("class_name") or "").strip()
                    if label:
                        searchable_parts.append(label)

        return "\n".join(searchable_parts)

    async def _process_html(self, file_path: str) -> str:
        """Process HTML document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                return soup.get_text()
        except Exception as e:
            logger.error(f"❌ HTML processing failed: {e}")
            return ""
    
    async def _process_eml(self, file_path: str) -> str:
        """Process EML email file with Unicode sanitization"""
        try:
            with open(file_path, 'rb') as file:
                msg = email.message_from_bytes(file.read())
            
            # Extract email metadata
            headers = []
            for header in ['From', 'To', 'Cc', 'Subject', 'Date']:
                value = msg.get(header)
                if value:
                    # Sanitize header value
                    clean_value = self._sanitize_unicode(str(value))
                    headers.append(f"{header}: {clean_value}")
            
            # Extract email body
            body_parts = []
            
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == 'text/plain':
                        try:
                            decoded_content = part.get_payload(decode=True).decode('utf-8', errors='replace')
                            clean_content = self._sanitize_unicode(decoded_content)
                            body_parts.append(clean_content)
                        except:
                            try:
                                decoded_content = part.get_payload(decode=True).decode('latin-1', errors='replace')
                                clean_content = self._sanitize_unicode(decoded_content)
                                body_parts.append(clean_content)
                            except:
                                continue
                    elif content_type == 'text/html':
                        try:
                            html_content = part.get_payload(decode=True).decode('utf-8', errors='replace')
                            soup = BeautifulSoup(html_content, 'html.parser')
                            text_content = soup.get_text()
                            clean_content = self._sanitize_unicode(text_content)
                            body_parts.append(clean_content)
                        except:
                            continue
            else:
                # Single part message
                try:
                    content = msg.get_payload(decode=True)
                    if isinstance(content, bytes):
                        content = content.decode('utf-8', errors='replace')
                    
                    if msg.get_content_type() == 'text/html':
                        soup = BeautifulSoup(content, 'html.parser')
                        text_content = soup.get_text()
                        clean_content = self._sanitize_unicode(text_content)
                        body_parts.append(clean_content)
                    else:
                        clean_content = self._sanitize_unicode(str(content))
                        body_parts.append(clean_content)
                except:
                    raw_payload = str(msg.get_payload())
                    clean_content = self._sanitize_unicode(raw_payload)
                    body_parts.append(clean_content)
            
            # Combine headers and body
            full_text = "\n".join(headers) + "\n\n" + "\n".join(body_parts)
            
            # Final sanitization of the complete text
            full_text = self._sanitize_unicode(full_text)
            
            logger.info(f"✅ EML processed: {len(headers)} headers, {len(body_parts)} body parts")
            return full_text
            
        except Exception as e:
            logger.error(f"❌ EML processing failed: {e}")
            return ""
    
    def _extract_subtitle_cues_text(self, content: str) -> tuple[str, int]:
        """
        Extract cue text from SRT or WebVTT. Finds the timestamp line (contains '-->')
        in each blank-line-separated block and keeps all following lines as subtitle text.
        Skips WEBVTT headers, NOTE/REGION blocks without timestamps, etc.
        """
        content = content.lstrip('\ufeff').replace('\r\n', '\n').replace('\r', '\n').strip()
        if not content:
            return "", 0
        blocks = re.split(r'\n\s*\n', content)
        subtitles: List[str] = []
        for block in blocks:
            block = block.strip()
            if not block:
                continue
            lines = block.split('\n')
            ts_idx = None
            for i, line in enumerate(lines):
                if '-->' in line:
                    ts_idx = i
                    break
            if ts_idx is None:
                continue
            subtitle_text = '\n'.join(lines[ts_idx + 1:]).strip()
            if subtitle_text:
                subtitles.append(subtitle_text)
        combined = '\n\n'.join(subtitles)
        return combined, len(subtitles)

    async def _process_subtitle_file(self, file_path: str, source_kind: str) -> str:
        """Process SRT or WebVTT: strip timing/metadata, return plain subtitle text for chunking."""
        label = 'SRT' if source_kind == 'srt' else 'WebVTT'
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            combined_text, n = self._extract_subtitle_cues_text(content)
            logger.info(f"✅ {label} processed: {n} cues extracted")
            return combined_text
        except UnicodeDecodeError:
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    combined_text, n = self._extract_subtitle_cues_text(content)
                    logger.info(f"✅ {label} processed with {encoding}: {n} cues extracted")
                    return combined_text
                except UnicodeDecodeError:
                    continue
            logger.error(f"❌ Could not decode {label} file: {file_path}")
            return ""
        except Exception as e:
            logger.error(f"❌ {label} processing failed: {e}")
            return ""
    
    async def _process_zip(self, file_path: str) -> str:
        """ZIP archives are expanded at upload time via ZipProcessorService; do not merge as one doc."""
        logger.info(
            "ZIP file reached _process_zip; ZIP contents are processed as individual "
            "documents via ZipProcessorService at upload time. Returning empty text. path=%s",
            file_path,
        )
        return ""
    
    async def _assess_quality(self, text: str, ocr_confidence: float) -> QualityMetrics:
        """Assess the quality of extracted text"""
        try:
            # Language detection confidence
            try:
                language = detect(text)
                lang_confidence = 0.9 if language == 'en' else 0.7
            except:
                lang_confidence = 0.5
            
            # Vocabulary coherence (ratio of dictionary words)
            words = text.split()
            if len(words) > 0:
                # Simple heuristic: check for reasonable word length distribution
                avg_word_length = sum(len(word) for word in words) / len(words)
                vocab_score = min(1.0, max(0.0, (avg_word_length - 2) / 10))
            else:
                vocab_score = 0.0
            
            # Text pattern analysis (punctuation, capitalization, etc.)
            if len(text) > 0:
                capital_ratio = sum(1 for c in text if c.isupper()) / len(text)
                punct_ratio = sum(1 for c in text if c in '.,!?;:') / len(text)
                pattern_score = min(1.0, (capital_ratio * 10 + punct_ratio * 20))
            else:
                pattern_score = 0.0
            
            # Overall score
            overall_score = (
                ocr_confidence * 0.3 +
                lang_confidence * 0.3 +
                vocab_score * 0.2 +
                pattern_score * 0.2
            )
            
            return QualityMetrics(
                ocr_confidence=ocr_confidence,
                language_confidence=lang_confidence,
                vocabulary_score=vocab_score,
                pattern_score=pattern_score,
                overall_score=overall_score
            )
            
        except Exception as e:
            logger.error(f"❌ Quality assessment failed: {e}")
            return QualityMetrics(
                ocr_confidence=ocr_confidence,
                language_confidence=0.5,
                vocabulary_score=0.5,
                pattern_score=0.5,
                overall_score=0.5
            )
    
    def _assign_page_range(
        self,
        chunk: Chunk,
        full_text: str,
        page_boundaries: list[tuple[int, int]],
    ) -> None:
        """Set chunk.metadata page_start and page_end from character offsets.
        page_boundaries is list of (page_num_1based, char_start). Mutates chunk in place.
        """
        if not page_boundaries:
            return
        content = chunk.content.strip()
        if not content:
            return
        start = full_text.find(content)
        if start < 0:
            return
        end = start + len(content)
        page_start = None
        page_end = None
        for i, (page_num, char_start) in enumerate(page_boundaries):
            if char_start <= start:
                page_start = page_num
            if char_start <= end:
                page_end = page_num
        if page_start is not None:
            chunk.metadata["page_start"] = page_start
            chunk.metadata["page_end"] = page_end if page_end is not None else page_start

    def _chunk_newspaper_content(self, text: str, file_path: str, document_id: str) -> List[Chunk]:
        """Chunk newspaper/magazine content by article boundaries (headline + body)."""
        chunks = []
        lines = text.splitlines()
        article_index = 0
        chunk_index = 0
        i = 0
        stem = Path(file_path).stem

        def is_headline(line: str) -> bool:
            line = line.strip()
            if not line or len(line) > 80:
                return False
            if line.isupper() and len(line) >= 10:
                return True
            words = line.split()
            if len(words) > 12:
                return False
            if words and words[0][0].isupper() and not line.endswith("."):
                return True
            return False

        while i < len(lines):
            line = lines[i]
            if not line.strip():
                i += 1
                continue
            if is_headline(line):
                headline = line.strip()
                body_lines = []
                i += 1
                while i < len(lines):
                    next_line = lines[i]
                    if next_line.strip() and is_headline(next_line):
                        break
                    body_lines.append(next_line)
                    i += 1
                body = "\n".join(body_lines).strip()
                if body:
                    paragraphs = [p.strip() for p in body.split("\n\n") if p.strip()]
                    for pi, para in enumerate(paragraphs):
                        if len(para) > 30:
                            chunks.append(self._create_chunk(
                                f"{stem}_art{article_index}_p{chunk_index}",
                                document_id,
                                para,
                                chunk_index,
                                "newspaper_paragraph",
                                {"article_headline": headline, "article_index": article_index, "paragraph_index": pi},
                            ))
                            chunk_index += 1
                else:
                    chunks.append(self._create_chunk(
                        f"{stem}_art{article_index}_h",
                        document_id,
                        headline,
                        chunk_index,
                        "newspaper_headline",
                        {"article_headline": headline, "article_index": article_index},
                    ))
                    chunk_index += 1
                article_index += 1
            else:
                acc = []
                while i < len(lines):
                    ln = lines[i]
                    if ln.strip() and is_headline(ln):
                        break
                    acc.append(ln)
                    i += 1
                body = "\n".join(acc).strip()
                if body and len(body) > 30:
                    chunks.append(self._create_chunk(
                        f"{stem}_lead_{chunk_index}",
                        document_id,
                        body,
                        chunk_index,
                        "newspaper_lead",
                        {"article_index": article_index},
                    ))
                    chunk_index += 1
                    article_index += 1

        return chunks

    async def _chunk_text(
        self,
        text: str,
        file_path: str,
        document_id: str,
        *,
        page_boundaries: list[tuple[int, int]] | None = None,
        doc_type_hint: str | None = None,
    ) -> List[Chunk]:
        """Universal chunking system that adapts to document type and structure.
        If page_boundaries is provided (e.g. from PDF), each chunk gets page_start/page_end in metadata.
        If doc_type_hint is 'newspaper', uses newspaper/magazine article-boundary chunking.
        """
        try:
            logger.info(f"Text length: {len(text)} characters")
            logger.info(f"First 200 chars: {text[:200]}...")

            if doc_type_hint == "newspaper":
                doc_structure = "newspaper"
            else:
                doc_structure = self._analyze_document_structure(text)
            logger.info(f"Document structure: {doc_structure}")

            chunks = []

            if doc_structure == "newspaper":
                chunks = self._chunk_newspaper_content(text, file_path, document_id)
            elif doc_structure == "book":
                chunks = self._chunk_book_content(text, file_path, document_id)
            elif doc_structure == "email":
                chunks = self._chunk_email_content(text, file_path, document_id)
            elif doc_structure == "academic_paper":
                chunks = self._chunk_academic_content(text, file_path, document_id)
            elif doc_structure == "article":
                chunks = self._chunk_article_content(text, file_path, document_id)
            else:
                # Default hierarchical chunking
                chunks = self._chunk_hierarchical(text, file_path, document_id)

            # Post-process: ensure optimal chunk sizes
            final_chunks = self._optimize_chunk_sizes(chunks, file_path, document_id)

            # Assign page range to each chunk when we have page boundaries (e.g. PDF)
            if page_boundaries and final_chunks:
                for ch in final_chunks:
                    self._assign_page_range(ch, text, page_boundaries)

            logger.info(f"Created {len(final_chunks)} total chunks using {doc_structure} strategy")
            return final_chunks

        except Exception as e:
            logger.error(f"Text chunking failed: {e}")
            return []
    
    def _analyze_document_structure(self, text: str) -> str:
        """Analyze document to determine its structure type"""
        text_lower = text.lower()
        lines = text.split('\n')
        
        # Email detection
        email_indicators = ['from:', 'to:', 'subject:', '@', 'sent:', 'cc:']
        if sum(1 for line in lines[:10] if any(ind in line.lower() for ind in email_indicators)) >= 2:
            return "email"
        
        # Book detection
        book_indicators = ['chapter', 'table of contents', 'preface', 'introduction', 'epilogue']
        if any(indicator in text_lower for indicator in book_indicators):
            # Check for chapter-like structure
            chapter_patterns = ['chapter ', 'part ', 'section ']
            if sum(1 for line in lines if any(pattern in line.lower() for pattern in chapter_patterns)) >= 2:
                return "book"
        
        # Academic paper detection
        academic_indicators = ['abstract', 'introduction', 'methodology', 'results', 'conclusion', 'references', 'bibliography']
        if sum(1 for indicator in academic_indicators if indicator in text_lower) >= 3:
            return "academic_paper"
        
        # Article detection (news, blog, etc.)
        if len(text) > 1000 and len(text) < 50000:  # Medium length
            paragraphs = [p for p in text.split('\n\n') if len(p.strip()) > 100]
            if len(paragraphs) >= 3:
                return "article"
        
        return "general"
    
    def _chunk_book_content(self, text: str, file_path: str, document_id: str) -> List[Chunk]:
        """Chunk book content by chapters, sections, and paragraphs"""
        chunks = []
        lines = text.split('\n')
        
        current_chapter = ""
        current_section = ""
        current_content = ""
        chunk_index = 0
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            line_type = self._classify_book_line(line)
            
            if line_type == "chapter_title":
                # Save previous content
                if current_content.strip():
                    chunks.append(self._create_chunk(
                        f"{document_id}_ch_{chunk_index}",
                        document_id,
                        current_content.strip(),
                        chunk_index,
                        "chapter_content",
                        {"chapter": current_chapter, "section": current_section}
                    ))
                    chunk_index += 1
                
                current_chapter = line
                current_section = ""
                current_content = ""
                
                # Chapter title as its own chunk
                chunks.append(self._create_chunk(
                    f"{Path(file_path).stem}_title_{chunk_index}",
                    document_id,
                    line,
                    chunk_index,
                    "chapter_title",
                    {"chapter": line}
                ))
                chunk_index += 1
                
            elif line_type == "section_title":
                # Save previous section content
                if current_content.strip():
                    chunks.append(self._create_chunk(
                        f"{Path(file_path).stem}_sec_{chunk_index}",
                        document_id,
                        current_content.strip(),
                        chunk_index,
                        "section_content",
                        {"chapter": current_chapter, "section": current_section}
                    ))
                    chunk_index += 1
                
                current_section = line
                current_content = line + "\n"
                
            else:
                current_content += line + "\n"
                
                # Split if content gets too long (increased threshold)
                if len(current_content) > 1500:
                    chunks.append(self._create_chunk(
                        f"{Path(file_path).stem}_para_{chunk_index}",
                        document_id,
                        current_content.strip(),
                        chunk_index,
                        "paragraph_group",
                        {"chapter": current_chapter, "section": current_section}
                    ))
                    chunk_index += 1
                    current_content = ""
        
        # Add final content
        if current_content.strip():
            chunks.append(self._create_chunk(
                f"{Path(file_path).stem}_final_{chunk_index}",
                document_id,
                current_content.strip(),
                chunk_index,
                "final_content",
                {"chapter": current_chapter, "section": current_section}
            ))
        
        return chunks
    
    def _chunk_email_content(self, text: str, file_path: str, document_id: str) -> List[Chunk]:
        """Chunk email content with aggressive size limits and token awareness"""
        chunks = []
        
        # Split the massive email collection into individual emails first
        email_files = text.split('\n==================================================\nFILE: ')
        
        chunk_index = 0
        
        for i, email_content in enumerate(email_files):
            if not email_content.strip():
                continue
            
            # Add back the file separator for non-first emails
            if i > 0:
                email_content = 'FILE: ' + email_content
            
            # Extract filename from the content if present
            lines = email_content.split('\n')
            file_name = "unknown"
            content_start = 0
            
            # Look for the filename in the first few lines
            for j, line in enumerate(lines[:3]):
                if line.startswith('FILE: ') or 'FILE:' in line:
                    file_name = line.replace('FILE:', '').strip()
                    content_start = j + 1
                    break
                elif line.startswith('=================================================='):
                    content_start = j + 1
                    break
            
            # Process the actual email content (skip file headers)
            email_body = '\n'.join(lines[content_start:])
            
            # Now chunk this individual email with strict token limits
            email_chunks = self._chunk_single_email(email_body, file_path, document_id, chunk_index, file_name)
            chunks.extend(email_chunks)
            chunk_index += len(email_chunks)
        
        return chunks
    
    def _chunk_single_email(self, email_text: str, file_path: str, document_id: str, start_index: int, file_name: str) -> List[Chunk]:
        """Chunk a single email with very strict token limits"""
        chunks = []
        lines = email_text.split('\n')
        
        current_content = ""
        current_type = "content"
        chunk_index = start_index
        max_tokens_per_chunk = 2500  # Reasonable limit that preserves semantic context
        
        # Process line by line with strict token monitoring
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Skip system messages and artifacts
            if self._is_system_message(line):
                continue
            
            # Determine content type
            line_type = self._classify_line_type(line)
            
            # Check if adding this line would exceed token limits
            test_content = current_content + "\n" + line if current_content else line
            estimated_tokens = self._estimate_chunk_tokens(test_content)
            
            # If adding this line would exceed limits, save current chunk and start new one
            if estimated_tokens > max_tokens_per_chunk and current_content.strip():
                chunks.append(self._create_chunk(
                    f"{Path(file_path).stem}_{current_type}_{chunk_index}",
                    document_id,
                    current_content.strip(),
                    chunk_index,
                    current_type,
                    {
                        "content_type": current_type,
                        "source_file": file_name,
                        "estimated_tokens": self._estimate_chunk_tokens(current_content)
                    }
                ))
                chunk_index += 1
                current_content = line
                current_type = line_type
            
            # If the content type changes significantly, start a new chunk
            elif current_type != line_type and current_content.strip() and len(current_content) > 100:
                chunks.append(self._create_chunk(
                    f"{Path(file_path).stem}_{current_type}_{chunk_index}",
                    document_id,
                    current_content.strip(),
                    chunk_index,
                    current_type,
                    {
                        "content_type": current_type,
                        "source_file": file_name,
                        "estimated_tokens": self._estimate_chunk_tokens(current_content)
                    }
                ))
                chunk_index += 1
                current_content = line
                current_type = line_type
            
            else:
                # Add line to current chunk
                current_content = test_content
                current_type = line_type
        
        # Add final chunk if exists
        if current_content.strip():
            chunks.append(self._create_chunk(
                f"{Path(file_path).stem}_{current_type}_{chunk_index}",
                document_id,
                current_content.strip(),
                chunk_index,
                current_type,
                {
                    "content_type": current_type,
                    "source_file": file_name,
                    "estimated_tokens": self._estimate_chunk_tokens(current_content)
                }
            ))
        
        return chunks
    
    def _is_system_message(self, line: str) -> bool:
        """Check if line is a system message to skip"""
        line_lower = line.lower()
        skip_patterns = [
            'microsoft office', 'prevented automatic', 'to help protect',
            'download of this picture', 'from the internet'
        ]
        return any(pattern in line_lower for pattern in skip_patterns)
    
    def _is_quote_line(self, line: str) -> bool:
        """Check if line looks like a quote"""
        # Must be reasonable length and end with punctuation
        if len(line) < 20 or len(line) > 200:
            return False
        
        # Must end with sentence-ending punctuation
        if not any(line.endswith(p) for p in ['.', '!', '?']):
            return False
        
        # Skip if it looks like email metadata
        line_lower = line.lower()
        if any(pattern in line_lower for pattern in ['from:', 'to:', 'subject:', '@', 'sent:']):
            return False
        
        return True
    
    def _chunk_academic_content(self, text: str, file_path: str, document_id: str) -> List[Chunk]:
        """Chunk academic papers by sections"""
        chunks = []
        sections = ['abstract', 'introduction', 'methodology', 'results', 'discussion', 'conclusion', 'references']
        
        current_section = ""
        current_content = ""
        chunk_index = 0
        
        paragraphs = text.split('\n\n')
        
        for paragraph in paragraphs:
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            
            # Check if this paragraph starts a new section
            para_lower = paragraph.lower()
            found_section = None
            for section in sections:
                if para_lower.startswith(section) or f"\n{section}" in para_lower:
                    found_section = section
                    break
            
            if found_section:
                # Save previous section
                if current_content.strip():
                    chunks.append(self._create_chunk(
                        f"{Path(file_path).stem}_{current_section}_{chunk_index}",
                        document_id,
                        current_content.strip(),
                        chunk_index,
                        f"academic_{current_section}",
                        {"section": current_section}
                    ))
                    chunk_index += 1
                
                current_section = found_section
                current_content = paragraph + "\n"
            else:
                current_content += paragraph + "\n"
                
                # Split long sections
                if len(current_content) > 1000:
                    chunks.append(self._create_chunk(
                        f"{Path(file_path).stem}_{current_section}_{chunk_index}",
                        document_id,
                        current_content.strip(),
                        chunk_index,
                        f"academic_{current_section}",
                        {"section": current_section}
                    ))
                    chunk_index += 1
                    current_content = ""
        
        # Add final content
        if current_content.strip():
            chunks.append(self._create_chunk(
                f"{Path(file_path).stem}_{current_section}_{chunk_index}",
                document_id,
                current_content.strip(),
                chunk_index,
                f"academic_{current_section}",
                {"section": current_section}
            ))
        
        return chunks
    
    def _chunk_article_content(self, text: str, file_path: str, document_id: str) -> List[Chunk]:
        """Chunk articles by paragraphs with context"""
        chunks = []
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        for i, paragraph in enumerate(paragraphs):
            if len(paragraph) > 50:  # Skip very short paragraphs
                chunk_type = "title" if i == 0 and len(paragraph) < 200 else "paragraph"
                
                chunks.append(self._create_chunk(
                    f"{Path(file_path).stem}_para_{i}",
                    document_id,
                    paragraph,
                    i,
                    chunk_type,
                    {"paragraph_number": i, "total_paragraphs": len(paragraphs)}
                ))
        
        return chunks
    
    def _chunk_hierarchical(self, text: str, file_path: str, document_id: str) -> List[Chunk]:
        """Default hierarchical chunking for general documents"""
        chunks = []
        
        # Try paragraph-based first
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        if len(paragraphs) > 1:
            for i, paragraph in enumerate(paragraphs):
                if len(paragraph) > 30:
                    chunks.append(self._create_chunk(
                        f"{Path(file_path).stem}_para_{i}",
                        document_id,
                        paragraph,
                        i,
                        "paragraph",
                        {}
                    ))
        else:
            # Fall back to fixed-size chunking
            chunk_size = 600
            overlap = 100
            
            for i in range(0, len(text), chunk_size - overlap):
                chunk_text = text[i:i + chunk_size]
                if chunk_text.strip():
                    chunks.append(self._create_chunk(
                        f"{Path(file_path).stem}_fixed_{i // (chunk_size - overlap)}",
                        document_id,
                        chunk_text.strip(),
                        i // (chunk_size - overlap),
                        "fixed_size",
                        {"start_pos": i, "end_pos": i + chunk_size}
                    ))
        
        return chunks
    
    def _classify_book_line(self, line: str) -> str:
        """Classify lines in book content"""
        line_lower = line.lower().strip()
        
        # Chapter titles
        if (line_lower.startswith('chapter ') or 
            line_lower.startswith('part ') or
            (len(line) < 100 and line.isupper())):
            return "chapter_title"
        
        # Section titles (often numbered or short)
        if (len(line) < 150 and 
            (line_lower.startswith(('1.', '2.', '3.', '4.', '5.')) or
             (len(line.split()) < 8 and not line.endswith('.')))):
            return "section_title"
        
        return "content"
    
    def _create_chunk(self, chunk_id: str, document_id: str, content: str, 
                     index: int, chunk_type: str, metadata: dict) -> Chunk:
        """Helper to create a chunk object with Unicode sanitization"""
        # Sanitize content to prevent Unicode encoding issues
        clean_content = self._sanitize_unicode(content)
        
        return Chunk(
            chunk_id=chunk_id,
            document_id=document_id,
            content=clean_content,
            chunk_index=index,
            quality_score=0.8,
            method="adaptive",
            metadata={
                "word_count": len(clean_content.split()),
                "char_count": len(clean_content),
                "chunk_type": chunk_type,
                **metadata
            }
        )
    
    def _optimize_chunk_sizes(self, chunks: List[Chunk], file_path: str, document_id: str) -> List[Chunk]:
        """Optimize chunk sizes for better embedding performance with very strict limits"""
        optimized_chunks = []
        
        for chunk in chunks:
            # Very aggressive size limits to prevent token overflow
            # Using conservative token estimation: 1 token per 2 characters
            estimated_tokens = self._estimate_chunk_tokens(chunk.content)
            
            if estimated_tokens > 4000:  # Balanced limit for better semantic context
                sub_chunks = self._split_large_chunk(chunk, document_id)
                optimized_chunks.extend(sub_chunks)
            elif len(chunk.content) < 50:  # Merge very small chunks
                # Try to merge with previous chunk if compatible
                if (optimized_chunks and 
                    self._estimate_chunk_tokens(optimized_chunks[-1].content) < 2000 and  # Safe merge threshold
                    optimized_chunks[-1].metadata.get('chunk_type') == chunk.metadata.get('chunk_type')):
                    
                    # Merge with previous chunk
                    prev_chunk = optimized_chunks[-1]
                    merged_content = prev_chunk.content + "\n" + chunk.content
                    
                    # Check if merged content is still within safe token limits
                    if self._estimate_chunk_tokens(merged_content) <= 3500:  # Reasonable merge limit
                        merged_chunk = Chunk(
                            chunk_id=f"{prev_chunk.chunk_id}_merged",
                            document_id=prev_chunk.document_id,
                            content=merged_content,
                            chunk_index=prev_chunk.chunk_index,
                            quality_score=prev_chunk.quality_score,
                            method="merged",
                            metadata={
                                "word_count": len(merged_content.split()),
                                "char_count": len(merged_content),
                                "chunk_type": prev_chunk.metadata.get('chunk_type'),
                                "merged_from": [prev_chunk.chunk_id, chunk.chunk_id],
                                "estimated_tokens": self._estimate_chunk_tokens(merged_content)
                            }
                        )
                        
                        optimized_chunks[-1] = merged_chunk
                    else:
                        # Don't merge if it would create too large a chunk
                        optimized_chunks.append(chunk)
                else:
                    optimized_chunks.append(chunk)
            else:
                # Add token estimate to metadata for monitoring
                chunk.metadata["estimated_tokens"] = estimated_tokens
                optimized_chunks.append(chunk)
        
        return optimized_chunks
    
    def _create_semantic_chunks(self, text: str) -> List[str]:
        """Create chunks based on semantic content boundaries"""
        lines = text.split('\n')
        chunks = []
        current_chunk = ""
        current_type = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            line_type = self._classify_line_type(line)
            
            # Start new chunk if content type changes or chunk gets too long
            if (current_type and line_type != current_type) or len(current_chunk) > 300:
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
                current_chunk = line
                current_type = line_type
            else:
                current_chunk += "\n" + line if current_chunk else line
                current_type = line_type
        
        # Add final chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _classify_line_type(self, line: str) -> str:
        """Classify what type of content a line contains"""
        line_lower = line.lower()
        
        # Email headers
        if any(pattern in line_lower for pattern in ['from:', 'to:', 'cc:', 'subject:', 'sent:', '@']):
            return 'email_header'
        
        # Quotes/sayings (usually end with punctuation and are standalone)
        if len(line) > 30 and any(line.endswith(p) for p in ['.', '!', '?']) and not line_lower.startswith(('from', 'to', 'cc', 'subject')):
            return 'quote'
        
        # Title/subject lines
        if 'bob hope' in line_lower and len(line) < 100:
            return 'title'
        
        # Technical/system messages
        if any(pattern in line_lower for pattern in ['microsoft office', 'prevented', 'download', 'protect']):
            return 'system_message'
        
        # Default content
        return 'content'
    
    def _classify_chunk_type(self, chunk: str) -> str:
        """Classify the overall type of a chunk"""
        chunk_lower = chunk.lower()
        
        if 'from:' in chunk_lower or 'to:' in chunk_lower:
            return 'email_metadata'
        elif 'bob hope' in chunk_lower and len(chunk) < 200:
            return 'title_or_subject'
        elif any(chunk.endswith(p) for p in ['.', '!', '?']) and len(chunk) > 30:
            return 'quote_or_saying'
        elif 'microsoft' in chunk_lower or 'office' in chunk_lower:
            return 'system_message'
        else:
            return 'general_content'
    
    def _split_large_chunk(self, chunk: Chunk, document_id: str) -> List[Chunk]:
        """Split a large chunk into smaller pieces with very strict token limits
        
        Note: document_id parameter is for consistency but chunk.document_id is already set correctly
        """
        content = chunk.content
        sentences = self._split_by_sentences(content)
        
        sub_chunks = []
        current_content = ""
        sub_index = 0
        max_sub_chunk_tokens = 2800  # Balanced token limit for good context
        
        for sentence in sentences:
            # Check if adding this sentence would exceed token limits
            test_content = current_content + " " + sentence if current_content else sentence
            estimated_tokens = self._estimate_chunk_tokens(test_content)
            
            if estimated_tokens > max_sub_chunk_tokens and current_content:
                # Save current chunk
                sub_chunk = Chunk(
                    chunk_id=f"{chunk.chunk_id}_sub_{sub_index}",
                    document_id=chunk.document_id,
                    content=current_content.strip(),
                    chunk_index=chunk.chunk_index * 1000 + sub_index,
                    quality_score=chunk.quality_score,
                    method="sub_chunk",
                    metadata={
                        "word_count": len(current_content.split()),
                        "char_count": len(current_content),
                        "estimated_tokens": self._estimate_chunk_tokens(current_content),
                        "parent_chunk": chunk.chunk_id,
                        "chunk_type": chunk.metadata.get('chunk_type', 'unknown'),
                        "parent_chunk_index": chunk.chunk_index,
                        "sub_index": sub_index
                    }
                )
                sub_chunks.append(sub_chunk)
                logger.debug(f"📝 Sub-chunk {sub_index}: {len(current_content)} chars, ~{self._estimate_chunk_tokens(current_content)} tokens")
                current_content = sentence
                sub_index += 1
            else:
                current_content = test_content
        
        # Add final sub-chunk
        if current_content.strip():
            sub_chunk = Chunk(
                chunk_id=f"{chunk.chunk_id}_sub_{sub_index}",
                document_id=chunk.document_id,
                content=current_content.strip(),
                chunk_index=chunk.chunk_index * 1000 + sub_index,
                quality_score=chunk.quality_score,
                method="sub_chunk",
                metadata={
                    "word_count": len(current_content.split()),
                    "char_count": len(current_content),
                    "estimated_tokens": self._estimate_chunk_tokens(current_content),
                    "parent_chunk": chunk.chunk_id,
                    "chunk_type": chunk.metadata.get('chunk_type', 'unknown'),
                    "parent_chunk_index": chunk.chunk_index,
                    "sub_index": sub_index
                }
            )
            sub_chunks.append(sub_chunk)
            logger.debug(f"📝 Final sub-chunk {sub_index}: {len(current_content)} chars, ~{self._estimate_chunk_tokens(current_content)} tokens")
        
        return sub_chunks if sub_chunks else [chunk]  # Return original if no splits made
    
    def _estimate_chunk_tokens(self, text: str) -> int:
        """Estimate tokens for chunking decisions - very conservative"""
        if not text:
            return 0
        
        # Very conservative estimation: 1 token per 2 characters
        # This is more conservative than typical 3-4 chars per token
        char_count = len(text)
        base_tokens = char_count // 2
        
        # Add overhead for punctuation and special characters
        punct_chars = sum(1 for c in text if not c.isalnum() and not c.isspace())
        punct_overhead = punct_chars * 0.5
        
        # Add overhead for word boundaries
        words = len(text.split())
        word_overhead = words * 0.3
        
        # Extra overhead for email content
        if any(pattern in text.lower() for pattern in ['from:', 'to:', 'subject:', '@']):
            email_overhead = char_count * 0.1
        else:
            email_overhead = 0
        
        total_estimated = int(base_tokens + punct_overhead + word_overhead + email_overhead)
        
        # Add 30% safety buffer
        return int(total_estimated * 1.3)
    
    def _clean_email_content(self, text: str) -> str:
        """Clean email headers and extract main content"""
        lines = text.split('\n')
        cleaned_lines = []
        
        # Skip email headers (lines that look like email metadata)
        skip_patterns = [
            'From:', 'To:', 'Cc:', 'Subject:', 'Sent:', 'Date:',
            '@', 'gcfl', 'pilbeam', 'Microsoft Office',
            'To help protect', 'prevented automatic download'
        ]
        
        for line in lines:
            line = line.strip()
            if line and not any(pattern.lower() in line.lower() for pattern in skip_patterns):
                # Skip very short lines that are likely metadata
                if len(line) > 10 or any(char in line for char in '.!?'):
                    cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines)
    
    def _split_by_sentences(self, text: str) -> List[str]:
        """Split text into sentences"""
        # Simple sentence splitting
        import re
        sentences = re.split(r'[.!?]+', text)
        return [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
    
    def _sanitize_unicode(self, text: str) -> str:
        """Sanitize text to remove problematic Unicode characters that cause JSON serialization issues"""
        if not text:
            return text
        
        try:
            # Remove surrogate characters and other problematic Unicode
            # that can't be encoded in UTF-8
            sanitized = ""
            for char in text:
                try:
                    # Try to encode the character to UTF-8
                    char.encode('utf-8')
                    sanitized += char
                except UnicodeEncodeError:
                    # Replace problematic characters with a safe placeholder
                    sanitized += '?'
            
            # Additional cleanup for common email encoding issues
            sanitized = sanitized.replace('\ufffd', '?')  # Replace replacement character
            sanitized = sanitized.replace('\u0000', '')   # Remove null characters
            
            # Remove or replace other problematic characters
            import re
            # Remove control characters except newlines, tabs, and carriage returns
            sanitized = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', sanitized)
            
            return sanitized
            
        except Exception as e:
            logger.warning(f"⚠️ Unicode sanitization failed: {e}")
            # Fallback: encode/decode to remove problematic characters
            try:
                return text.encode('utf-8', errors='replace').decode('utf-8')
            except:
                return str(text)
    
    def _map_entities_to_chunks(self, entities: List[Entity], chunks: List[Chunk]) -> List[Entity]:
        """Map each entity to the chunk(s) where its name appears (substring match, case-insensitive)."""
        if not chunks:
            return entities
        result = []
        for entity in entities:
            name_lower = entity.name.lower()
            chunk_ids = [
                c.chunk_id for c in chunks
                if name_lower in (c.content or "").lower()
            ]
            first_chunk = chunk_ids[0] if chunk_ids else ""
            meta = dict(entity.metadata or {})
            meta["chunk_ids"] = chunk_ids
            result.append(
                entity.model_copy(
                    update={
                        "source_chunk": first_chunk,
                        "metadata": meta,
                    }
                )
            )
        return result

    async def _extract_entities(self, text: str, chunks: List[Chunk]) -> List[Entity]:
        """Extract entities via document-service (spaCy + patterns + dedup all in document-service)."""
        try:
            logger.info("Starting entity extraction...")
            if self.document_service_client and getattr(
                self.document_service_client, "_initialized", False
            ):
                try:
                    entities = await self.document_service_client.extract_entities(text)
                    logger.info("Document service returned %s entities", len(entities))
                    entities = self._map_entities_to_chunks(entities, chunks)
                    return entities
                except Exception as e:
                    logger.warning("Document service entity extraction failed: %s", e)
            else:
                logger.warning("Document service not available for entity extraction")
            return []
        except Exception as e:
            logger.error("Entity extraction failed: %s", e)
            return []

    async def process_text_content(self, content: str, document_id: str, metadata: Dict[str, Any] = None) -> List[Chunk]:
        """Process text content directly and return chunks"""
        try:
            logger.info(f"🔄 Processing text content for document: {document_id}")
            
            # Clean and normalize text
            content = self._sanitize_unicode(content)
            
            # Assess quality
            quality_metrics = await self._assess_quality(content, 1.0)
            
            # Create chunks
            chunks = await self._chunk_text(content, f"{document_id}.txt", document_id)
            
            # Add metadata to chunks
            for chunk in chunks:
                chunk.metadata = chunk.metadata or {}
                chunk.metadata.update(metadata or {})
                chunk.metadata["quality_metrics"] = quality_metrics.dict()
                chunk.metadata["processing_method"] = "text_content"
                chunk.metadata["document_id"] = document_id
            
            logger.info(f"✅ Processed text content: {len(chunks)} chunks for {document_id}")
            return chunks
            
        except Exception as e:
            logger.error(f"❌ Failed to process text content for {document_id}: {e}")
            return []


# Global document processor instance for lazy loading
_document_processor_instance = None


async def get_document_processor() -> DocumentProcessor:
    """Get or create a global document processor instance"""
    global _document_processor_instance
    
    if _document_processor_instance is None:
        logger.info("🔄 Creating global document processor instance...")
        _document_processor_instance = DocumentProcessor()
        await _document_processor_instance.initialize()
        logger.info("✅ Global document processor initialized")
    
    return _document_processor_instance
